from pptx import Presentation
from .base import DocumentParser, DocumentPage, TextBlock

def pptx_table_to_markdown(table) -> str:
    """Convert python-pptx table to markdown format."""
    rows_data = []
    for row in table.rows:
        row_cells = []
        for cell in row.cells:
            row_cells.append(cell.text.strip().replace("\n", " "))
        rows_data.append(row_cells)
        
    if not rows_data:
        return ""
        
    md = []
    header = rows_data[0]
    md.append("| " + " | ".join(header) + " |")
    md.append("| " + " | ".join(["---"] * len(header)) + " |")
    
    for row in rows_data[1:]:
        md.append("| " + " | ".join(row) + " |")
        
    return "\n".join(md)

class PptxParser(DocumentParser):
    def parse(self, path: str, include_pages: str | None = None) -> list[DocumentPage]:
        prs = Presentation(path)
        
        pages_out = []
        for slide_idx, slide in enumerate(prs.slides):
            page_num = slide_idx + 1
            text_blocks = []
            tables_out = []
            
            # Check for slide title shape
            title_shape = None
            try:
                if slide.shapes.title:
                    title_shape = slide.shapes.title
            except AttributeError:
                pass
            
            for shape in slide.shapes:
                # 1. Text extraction
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if not text:
                            continue
                            
                        # Extract formatting
                        sizes = []
                        is_bold = False
                        for run in paragraph.runs:
                            if run.font.size:
                                sizes.append(run.font.size.pt)
                            if run.font.bold:
                                is_bold = True
                                
                        avg_size = sum(sizes) / len(sizes) if sizes else 14.0
                        
                        # Apply special heuristics for titles
                        if shape == title_shape or (shape.name and shape.name.startswith("Title")):
                            is_bold = True
                            if not sizes:
                                avg_size = 24.0
                                
                        text_blocks.append(TextBlock(
                            text=text,
                            font_size=avg_size,
                            is_bold=is_bold
                        ))
                        
                # 2. Table extraction
                if shape.has_table:
                    md_table = pptx_table_to_markdown(shape.table)
                    if md_table:
                        tables_out.append(md_table)
            
            pages_out.append(DocumentPage(
                page_number=page_num,
                text_blocks=text_blocks,
                tables=tables_out,
                source_path=path
            ))
            
        return pages_out
